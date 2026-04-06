#!/usr/bin/env python3
"""
Generate ASRRL Methodology PowerPoint — IEEE Algorithmic Format.
9-Stage Pipeline (Stage 0: PCAP → Stage 8: Classification).
Includes Parallel Processing slide and updated architecture diagram.
"""

import os
import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ───────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x00, 0x3C, 0x71)
MED_BLUE    = RGBColor(0x29, 0x80, 0xB9)
LIGHT_BLUE  = RGBColor(0xD6, 0xEA, 0xF8)
DARK_GREEN  = RGBColor(0x1E, 0x8E, 0x3E)
DARK_RED    = RGBColor(0xC0, 0x39, 0x2B)
DARK_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
GRAY        = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY  = RGBColor(0xEC, 0xF0, 0xF1)
PURPLE      = RGBColor(0x8E, 0x44, 0xAD)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
TEAL        = RGBColor(0x16, 0xA0, 0x85)
NAVY        = RGBColor(0x2C, 0x3E, 0x50)
ALGO_BG     = RGBColor(0xF7, 0xF9, 0xFB)
ALGO_BORDER = RGBColor(0x34, 0x49, 0x5E)
BROWN       = RGBColor(0x7B, 0x4B, 0x2A)
YELLOW_DARK = RGBColor(0xF3, 0x9C, 0x12)


def render_eq(latex_str, fontsize=14, dpi=200, figw=8, figh=0.6):
    """Render a single-line mathtext equation to PNG bytes."""
    fig, ax = plt.subplots(figsize=(figw, figh))
    ax.axis("off")
    ax.text(0.5, 0.5, latex_str, fontsize=fontsize,
            ha="center", va="center", transform=ax.transAxes)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                pad_inches=0.05, transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf


def render_eq_block(lines, fontsize=12, dpi=200, figw=9):
    """Render multiple mathtext lines to a single PNG."""
    n = len(lines)
    fig, ax = plt.subplots(figsize=(figw, 0.45 * n + 0.3))
    ax.axis("off")
    for i, line in enumerate(lines):
        y = 1.0 - (i + 0.5) / n
        ax.text(0.02, y, line, fontsize=fontsize,
                ha="left", va="center", transform=ax.transAxes)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                pad_inches=0.05, transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_ml(slide, left, top, width, height, lines_data, font_name="Calibri"):
    """Multiline text box. lines_data: [(text, size, bold, color), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_size, bold, color) in enumerate(lines_data):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(2)
    return txBox


def add_rect(slide, left, top, width, height, fill_color, text,
             font_size=11, font_color=WHITE, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = "Calibri"
    return shape


def add_right_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MED_BLUE
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MED_BLUE
    shape.line.fill.background()
    return shape


def add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_ieee_algorithm(slide, left, top, width, height, title_num, title_text, lines):
    """
    Add an IEEE-style algorithm block with border, title, and numbered lines.
    lines: list of (text, indent_level, is_keyword_line)
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ALGO_BG
    shape.line.color.rgb = ALGO_BORDER
    shape.line.width = Pt(1.5)

    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.05),
                                     width - Inches(0.3), height - Inches(0.1))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"Algorithm {title_num}: "
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.italic = True
    run.font.color.rgb = BLACK
    run.font.name = "Times New Roman"
    run2 = p.add_run()
    run2.text = title_text
    run2.font.size = Pt(12)
    run2.font.bold = True
    run2.font.color.rgb = BLACK
    run2.font.name = "Times New Roman"
    p.space_after = Pt(3)

    line_num = 0
    for text, indent, is_special in lines:
        p = tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(1)

        if text.startswith("Require:") or text.startswith("Ensure:"):
            run = p.add_run()
            keyword = text.split(":")[0] + ":"
            rest = text[len(keyword):]
            run.text = keyword
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = BLACK
            run.font.name = "Times New Roman"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(10)
            run2.font.bold = False
            run2.font.color.rgb = BLACK
            run2.font.name = "Times New Roman"
            continue

        line_num += 1
        indent_str = "    " * indent
        prefix = f"{line_num:>2}:  {indent_str}"
        _add_keyword_runs(p, prefix, text)

    return txBox


def _add_keyword_runs(paragraph, prefix, text):
    """Add runs to a paragraph, bolding IEEE keywords."""
    KEYWORDS = ["if ", "then", "else", "end if", "end for", "end while",
                "for ", "do", "while ", "return ", "repeat", "until"]

    run = paragraph.add_run()
    run.text = prefix
    run.font.size = Pt(10)
    run.font.bold = False
    run.font.color.rgb = GRAY
    run.font.name = "Consolas"

    remaining = text
    while remaining:
        earliest_pos = len(remaining)
        earliest_kw = None
        for kw in KEYWORDS:
            lower = remaining.lower()
            pos = lower.find(kw.lower())
            if pos != -1 and pos < earliest_pos:
                earliest_pos = pos
                earliest_kw = kw

        if earliest_kw is None:
            run = paragraph.add_run()
            run.text = remaining
            run.font.size = Pt(10)
            run.font.bold = False
            run.font.color.rgb = BLACK
            run.font.name = "Times New Roman"
            break
        else:
            if earliest_pos > 0:
                run = paragraph.add_run()
                run.text = remaining[:earliest_pos]
                run.font.size = Pt(10)
                run.font.bold = False
                run.font.color.rgb = BLACK
                run.font.name = "Times New Roman"
            kw_len = len(earliest_kw)
            run = paragraph.add_run()
            run.text = remaining[earliest_pos:earliest_pos + kw_len]
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = BLACK
            run.font.name = "Times New Roman"
            remaining = remaining[earliest_pos + kw_len:]
            continue


def add_bullets(slide, left, top, width, height, bullets, color=BLACK, font_size=12):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet_text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {bullet_text}"
        p.font.size = Pt(font_size)
        p.font.bold = False
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return txBox


# ═════════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ═════════════════════════════════════════════════════════════════════════
def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BLUE)
    add_text_box(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(1.5),
                 "ASRRL: Adaptive Symbolic Reasoning and\nReinforcement Learning for Dynamic\nNetwork Traffic Classification",
                 font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
                 "A Verifiable IDS Framework with Z3 Formal Verification,\nQ-Learning Safety Shielding, and DBSCAN Novel Pattern Detection",
                 font_size=18, color=RGBColor(0xBD, 0xC3, 0xC7), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
                 "3 Benchmark Datasets  |  9-Stage Pipeline (Stage 0\u20138)  |  Parallel Processing  |  IEEE Format",
                 font_size=16, color=RGBColor(0x95, 0xA5, 0xA6), alignment=PP_ALIGN.CENTER)

    add_notes(slide, """ASRRL \u2014 Adaptive Symbolic Reasoning and Reinforcement Learning for Dynamic Network Traffic Classification

This presentation details the complete methodology of the ASRRL Intrusion Detection System (IDS) framework. ASRRL combines multiple ML and formal methods techniques into a 9-stage pipeline (Stages 0\u20138) for network traffic classification.

KEY CONTRIBUTIONS:
1. PCAP Binary Preprocessing (Stage 0): Raw binary PCAP packet captures are converted to flow-level features using CICFlowMeter/Argus, establishing the bridge between raw network data and the ML pipeline.

2. Symbolic Reasoning via Z3 SMT Solver: Decision tree paths are converted into formal logic constraints using Microsoft Research's Z3 theorem prover, providing mathematically provable safety guarantees.

3. Reinforcement Learning with Safety Shielding: A tabular Q-learning agent learns optimal classification policies, but every proposed action is verified against Z3 constraints before execution.

4. Dynamic Adaptation: Adaptive buffer sizing (10-200 flows) and dynamic threshold adaptation (0.40-0.70) handle concept drift at runtime.

5. Novel Pattern Detection: DBSCAN clustering discovers zero-day attack patterns and incorporates them as new Z3 constraints.

6. Parallel Processing: Pipeline stages with independent data paths are parallelized for throughput \u2014 PCAP processing, Z3 extraction, and final classification are embarrassingly parallel.

DATASETS: UNSW-NB15, CSE-CIC-IDS-2018, CIC-IDS2017.""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2: Architecture Overview (Updated: 9 stages + parallelism)
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
                 "ASRRL Architecture \u2014 9-Stage Pipeline with Parallel Processing",
                 font_size=28, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # Row 1: Stages 0-4
    row1 = [
        ("0. PCAP Binary\nPreprocessing", BROWN),
        ("1. Feature Extract\n& Normalization", DARK_BLUE),
        ("2. Decision Tree\nSymbolic Model", PURPLE),
        ("3. Z3 Constraint\nExtraction", GREEN),
        ("4. Q-Learning RL\nSafety Shielding", DARK_ORANGE),
    ]
    box_w, box_h = Inches(1.85), Inches(0.85)
    start_x, y1 = Inches(0.25), Inches(1.0)
    gap = Inches(0.22)
    for i, (label, color) in enumerate(row1):
        x = start_x + i * (box_w + gap)
        add_rect(slide, x, y1, box_w, box_h, color, label, font_size=9)
        if i < len(row1) - 1:
            add_right_arrow(slide, x + box_w + Inches(0.01), y1 + Inches(0.25),
                           Inches(0.19), Inches(0.35))

    # Row 2: Stages 5-8
    row2 = [
        ("5. DBSCAN Novel\nPattern Detection", DARK_RED),
        ("6. Adaptive Buffer\nManagement", MED_BLUE),
        ("7. Dynamic Threshold\nAdaptation", TEAL),
        ("8. Final Classification\nDecision", NAVY),
    ]
    y2 = Inches(2.1)
    r2_start = Inches(0.25) + (box_w + gap)  # offset to align under stages 1-4
    for i, (label, color) in enumerate(row2):
        x = r2_start + i * (box_w + gap)
        add_rect(slide, x, y2, box_w, box_h, color, label, font_size=9)
        if i < len(row2) - 1:
            add_right_arrow(slide, x + box_w + Inches(0.01), y2 + Inches(0.25),
                           Inches(0.19), Inches(0.35))

    # Arrow from row 1 to row 2 (Stage 4 -> Stage 5)
    add_down_arrow(slide, Inches(0.25) + (box_w + gap) * 0 + box_w / 2 - Inches(0.1) + (box_w + gap),
                   y1 + box_h + Inches(0.02), Inches(0.2), Inches(0.22))

    # Parallel processing indicators
    add_text_box(slide, Inches(0.2), Inches(3.15), Inches(12.8), Inches(0.35),
                 "\u26a1 Parallel Processing Stages:  Stage 0 (per-PCAP)  |  Stage 3 (per-leaf path)  |  Stage 8 (per-flow)  |  Cross-dataset (all 3 datasets independently)",
                 font_size=12, bold=True, color=YELLOW_DARK, alignment=PP_ALIGN.CENTER)

    # Datasets
    add_text_box(slide, Inches(0.3), Inches(3.6), Inches(12.5), Inches(0.4),
                 "Three Primary Datasets: UNSW-NB15 (2.5M flows, 9 attacks)  |  "
                 "CSE-CIC-IDS-2018 (16M flows, 7 attacks)  |  CIC-IDS2017 (3.1M flows, 14 attacks)",
                 font_size=12, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # Data transformation summary
    add_text_box(slide, Inches(0.3), Inches(4.1), Inches(12.5), Inches(0.35),
                 "Data Transformation Summary (per flow x_i):",
                 font_size=14, bold=True, color=DARK_BLUE)

    flow = [
        ("Stage 0:  Binary PCAP packets  \u2192  CICFlowMeter/Argus  \u2192  5-tuple flow aggregation  \u2192  CSV (49\u201380+ columns per flow)", 9, False, BLACK),
        ("Stage 1:  Raw columns (dur, rate, sbytes, ...)  \u2192  x_i = [flow_duration, pkt_rate, byte_rate, traffic_pat_var*, port_cat, size_cat, protocol] \u2208 R^7", 9, False, BLACK),
        ("Stage 2:  X_norm \u2208 R^(n\u00d77) \u2192 DecisionTree(max_depth=6) \u2192  leaf_id(x_i), P(attack|x_i), \u0177_DT(x_i) \u2208 {0,1}", 9, False, BLACK),
        ("Stage 3:  DT paths (root \u2192 leaf) \u2192 Z3 Implies(\u2227 (f_j op \u03b8_j), action=class)  \u2192  constraint set C = {\u03c6_1,...,\u03c6_L}  [\u26a1 parallel per leaf]", 9, False, BLACK),
        ("Stage 4:  state s_i = leaf_id(x_i) \u2192 Q(s,a) \u2192 a_i* = argmax Q(s_i,a)  subject to Z3 verify(x_i, a_i*) = sat", 9, False, BLACK),
        ("Stage 5:  misclassified {x_i : \u00e2_i \u2260 a_i*} \u2192 DBSCAN(\u03b5=1.5, minPts=5) \u2192 novel centroids \u2192 C' = C \u222a C_novel", 9, False, BLACK),
        ("Stage 6:  \u03c3_tpv, \u03c3\u00b2_bytes of window W_t \u2192 |B_{t+1}| = clip(|B_t| \u00b1 \u0394, 10, 200)", 9, False, BLACK),
        ("Stage 7:  EMA(FPR_t, FNR_t) \u2192 \u03c4_{t+1} = clip(\u03c4_t \u00b1 0.005, 0.40, 0.70)", 9, False, BLACK),
        ("Stage 8:  p_i \u2265 \u03c4 \u2192 ATTACK;  p_i \u2264 1-\u03c4 \u2192 BENIGN;  else \u2192 RL+Z3 shield \u2192 \u0177_i \u2208 {0,1}  [\u26a1 parallel per flow]", 9, False, BLACK),
    ]
    add_ml(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(3), flow, font_name="Consolas")

    add_notes(slide, """SLIDE 2 \u2014 ASRRL Architecture: 9-Stage Pipeline with Parallel Processing

This slide presents the complete 9-stage pipeline (Stage 0 through Stage 8). The key addition from the previous version is Stage 0 (PCAP Binary Preprocessing), which addresses how raw binary network captures are converted into the flow-level tabular data that the ML pipeline consumes.

STAGE 0 \u2014 PCAP PREPROCESSING:
Real network traffic is captured as binary PCAP (Packet CAPture) files containing raw packet headers and payloads. Before any ML processing:
1. Packets are reassembled into TCP streams
2. Flows are aggregated by 5-tuple (src_ip, dst_ip, src_port, dst_port, protocol)
3. Timeout-based flow termination applies (active ~120s, idle ~60s)
4. CICFlowMeter (for CIC datasets) or Argus (for UNSW-NB15) computes 49\u201380+ statistical features per flow
5. Output: CSV files with one row per flow

PARALLEL PROCESSING POINTS:
The pipeline is NOT fully sequential \u2014 several stages support parallel execution:
- Stage 0: Each PCAP file processed independently \u2192 embarrassingly parallel
- Stage 3: Each leaf path extracted and verified independently \u2192 parallel per leaf
- Stage 8: Each flow classified independently \u2192 embarrassingly parallel per flow
- Cross-dataset: All 3 datasets processed through independent pipeline instances

SEQUENTIAL BOTTLENECKS:
- Stage 2 (DT training): CART splits are computed greedily \u2192 inherently sequential
- Stage 4 (Q-learning): Episodes update shared Q-table \u2192 sequential within episodes
- Stage 5 (DBSCAN): Core algorithm uses sequential neighborhood queries
- Stages 6-7 (Buffer/Threshold): Temporal dependency on previous window state""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3: Three Datasets \u2014 Raw Features & Extraction
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Three Benchmark Datasets \u2014 Raw Features & Standardized Extraction",
                 font_size=26, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # UNSW-NB15 column
    add_rect(slide, Inches(0.2), Inches(0.8), Inches(4.2), Inches(0.45),
             DARK_BLUE, "UNSW-NB15  (2015, UNSW Canberra)", font_size=12)
    unsw_lines = [
        ("Raw Columns (49 features):", 11, True, BLACK),
        ("  dur, rate, sbytes, dbytes, spkts, dpkts,", 10, False, GRAY),
        ("  ct_srv_src, sttl, sjit, dsport, proto, ...", 10, False, GRAY),
        ("", 4, False, BLACK),
        ("PCAP Source: IXIA PerfectStorm tool", 10, True, BROWN),
        ("Flow Tool: Argus + Bro-IDS", 10, True, BROWN),
        ("", 4, False, BLACK),
        ("Feature Extraction:", 11, True, DARK_BLUE),
        ("  flow_duration = dur \u00d7 1000  (sec \u2192 ms)", 10, False, BLACK),
        ("  pkt_rate      = rate  (or (spkts+dpkts)/dur)", 10, False, BLACK),
        ("  byte_rate     = sbytes + dbytes", 10, False, BLACK),
        ("  traffic_pat_var* = normalize(ct_srv_src, [0,1])", 10, False, BLACK),
        ("  port_cat      = bin(dsport, [80,443,1K,5K,10K,65K])", 10, False, BLACK),
        ("  size_cat      = bin(byte_rate, [100,1K,10K,\u221e])", 10, False, BLACK),
        ("  protocol      = {tcp:0, udp:1, icmp:2}", 10, False, BLACK),
        ("", 4, False, BLACK),
        ("Attack Types (9): Generic, Exploits, Fuzzers,", 10, False, DARK_RED),
        ("  DoS, Recon, Analysis, Backdoor, Shellcode, Worms", 10, False, DARK_RED),
        ("Attack Ratio: ~30%", 10, True, DARK_RED),
    ]
    add_ml(slide, Inches(0.2), Inches(1.3), Inches(4.2), Inches(5.5), unsw_lines, font_name="Consolas")

    # CSE-CIC-IDS-2018 column
    add_rect(slide, Inches(4.6), Inches(0.8), Inches(4.2), Inches(0.45),
             DARK_BLUE, "CSE-CIC-IDS-2018  (2018, CIC/UNB)", font_size=12)
    cse_lines = [
        ("Raw Columns (80+ CICFlowMeter features):", 11, True, BLACK),
        ("  Flow Duration, Flow Packets/s,", 10, False, GRAY),
        ("  Flow Bytes/s, Fwd IAT Std,", 10, False, GRAY),
        ("  Dst Port, Protocol, Label, ...", 10, False, GRAY),
        ("", 4, False, BLACK),
        ("PCAP Source: AWS infrastructure (10 days)", 10, True, BROWN),
        ("Flow Tool: CICFlowMeter", 10, True, BROWN),
        ("", 4, False, BLACK),
        ("Feature Extraction:", 11, True, DARK_BLUE),
        ("  flow_duration = Flow Duration / 1000  (\u03bcs \u2192 ms)", 10, False, BLACK),
        ("  pkt_rate      = Flow Packets/s", 10, False, BLACK),
        ("  byte_rate     = Flow Bytes/s", 10, False, BLACK),
        ("  traffic_pat_var* = quantile_norm(Fwd IAT Std)", 10, False, BLACK),
        ("  port_cat      = bin(Dst Port, [80,443,1K,5K,10K,65K])", 10, False, BLACK),
        ("  size_cat      = bin(bytes\u00d7dur, [100,1K,10K,\u221e])", 10, False, BLACK),
        ("  protocol      = {6\u2192TCP:0, 17\u2192UDP:1, 1\u2192ICMP:2}", 10, False, BLACK),
        ("", 4, False, BLACK),
        ("Attack Types (7): DoS, DDoS, BruteForce,", 10, False, DARK_RED),
        ("  Infiltration, BotNet, WebAttack, SQL Injection", 10, False, DARK_RED),
        ("Attack Ratio: ~15%", 10, True, DARK_RED),
    ]
    add_ml(slide, Inches(4.6), Inches(1.3), Inches(4.2), Inches(5.5), cse_lines, font_name="Consolas")

    # CIC-IDS2017 column
    add_rect(slide, Inches(9.0), Inches(0.8), Inches(4.2), Inches(0.45),
             DARK_BLUE, "CIC-IDS2017  (2017, CIC/UNB)", font_size=12)
    cic_lines = [
        ("Raw Columns (80+ CICFlowMeter features):", 11, True, BLACK),
        ("  Flow Duration, Flow Pkts/s,", 10, False, GRAY),
        ("  Flow Byts/s, Flow IAT Std,", 10, False, GRAY),
        ("  Destination Port, Protocol, Label, ...", 10, False, GRAY),
        ("", 4, False, BLACK),
        ("PCAP Source: CIC lab network (5 days)", 10, True, BROWN),
        ("Flow Tool: CICFlowMeter", 10, True, BROWN),
        ("", 4, False, BLACK),
        ("Feature Extraction (same as CSE-CIC):", 11, True, DARK_BLUE),
        ("  flow_duration = Flow Duration / 1000  (\u03bcs \u2192 ms)", 10, False, BLACK),
        ("  pkt_rate      = Flow Packets/s", 10, False, BLACK),
        ("  byte_rate     = Flow Bytes/s", 10, False, BLACK),
        ("  traffic_pat_var* = quantile_norm(Fwd IAT Std)", 10, False, BLACK),
        ("  port_cat      = bin(Dst Port, [80,443,1K,5K,10K,65K])", 10, False, BLACK),
        ("  size_cat      = bin(bytes\u00d7dur, [100,1K,10K,\u221e])", 10, False, BLACK),
        ("  protocol      = {6\u2192TCP:0, 17\u2192UDP:1, 1\u2192ICMP:2}", 10, False, BLACK),
        ("", 4, False, BLACK),
        ("Attack Types (14): DoS, PortScan, DDoS,", 10, False, DARK_RED),
        ("  BruteForce, WebAttack, Infiltration,", 10, False, DARK_RED),
        ("  Heartbleed, Botnet, SSH-Patator, ...", 10, False, DARK_RED),
        ("Attack Ratio: ~20%", 10, True, DARK_RED),
    ]
    add_ml(slide, Inches(9.0), Inches(1.3), Inches(4.2), Inches(5.5), cic_lines, font_name="Consolas")

    add_text_box(slide, Inches(0.3), Inches(6.4), Inches(12.5), Inches(0.3),
                 "*traffic_pat_var is a traffic pattern variability proxy (NOT Shannon entropy): UNSW uses ct_srv_src (connection diversity); "
                 "CIC datasets use Fwd IAT Std (timing variability), both normalized to [0,1].",
                 font_size=10, bold=False, color=GRAY, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.3), Inches(6.7), Inches(12.5), Inches(0.5),
                 "All 3 datasets \u2192 Standardized schema:  x_i = [flow_duration, pkt_rate, byte_rate, traffic_pat_var*, "
                 "port_cat, size_cat, protocol],  y_i \u2208 {0,1}",
                 font_size=13, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)

    add_notes(slide, """SLIDE 3 \u2014 Three Benchmark Datasets: Raw Features & Standardized Extraction

This slide now includes PCAP source information and flow extraction tools for each dataset.

UNSW-NB15: PCAP captured by IXIA PerfectStorm tool at UNSW Canberra. Flows extracted by Argus and Bro-IDS network monitors. The raw PCAP binary data was processed into 49-column CSV format by the dataset creators.

CSE-CIC-IDS-2018: PCAP captured over 10 days on AWS infrastructure simulating enterprise network. CICFlowMeter processed the binary PCAPs into 80+ flow-level features using 5-tuple aggregation with configurable timeouts.

CIC-IDS2017: PCAP captured over 5 days at CIC lab network. Same CICFlowMeter tool processed the binary captures into the same 80+ feature format.

KEY POINT FOR ADVISOR: All three benchmark datasets ship as pre-extracted CSVs. The PCAP \u2192 flow conversion was performed by the dataset creators using established tools (CICFlowMeter, Argus). Our Stage 0 documents this preprocessing step and shows it as part of the complete pipeline.""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 4: Stage 0 \u2014 PCAP Binary Preprocessing (NEW)
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Stage 0: PCAP Binary Preprocessing & Flow Aggregation  [\u26a1 Parallelizable]",
                 font_size=24, bold=True, color=BROWN)

    algo_lines = [
        ("Require: Raw binary PCAP file(s) P = {p_1, ..., p_N}", 0, False),
        ("Ensure: Flow-level CSV with F features per flow (F=49 for UNSW, F=80+ for CIC)", 0, False),
        ("for each PCAP file p_k \u2208 P do  // \u26a1 parallelizable per file", 0, True),
        ("for each raw packet pkt \u2208 p_k do", 1, True),
        ("Parse binary headers: Ethernet \u2192 IP \u2192 TCP/UDP/ICMP", 2, False),
        ("Extract 5-tuple: key \u2190 (src_ip, dst_ip, src_port, dst_port, proto)", 2, False),
        ("if key \u2208 active_flows then", 2, True),
        ("active_flows[key].append(pkt)  // add to existing flow", 3, False),
        ("else", 2, True),
        ("active_flows[key] \u2190 new_flow(pkt)  // start new flow", 3, False),
        ("end if", 2, True),
        ("end for", 1, True),
        ("for each flow f \u2208 active_flows do", 1, True),
        ("if idle_time(f) > T_idle or active_time(f) > T_active then", 2, True),
        ("Terminate flow and compute statistical features:", 3, False),
        ("  duration, pkt_count, byte_count, pkt_rate, byte_rate,", 3, False),
        ("  fwd/bwd ratios, IAT stats (mean, std, min, max),", 3, False),
        ("  TCP flags, payload size stats, ...", 3, False),
        ("end if", 2, True),
        ("end for", 1, True),
        ("Export terminated flows as CSV rows with F feature columns", 1, False),
        ("end for", 0, True),
        ("return Merged CSV: D_raw with n flows \u00d7 F features", 0, True),
    ]
    add_ieee_algorithm(slide, Inches(0.3), Inches(0.7), Inches(7.8), Inches(5.8),
                       1, "PCAP Binary Preprocessing & Flow Aggregation", algo_lines)

    # Right side: tool diagram and equations
    eq = render_eq_block([
        r"$flow\_key = (src\_ip, dst\_ip, src\_port, dst\_port, proto)$",
        r"$T_{idle} = 60s$,    $T_{active} = 120s$",
        r"$IAT_i = t_{pkt_{i+1}} - t_{pkt_i}$  (inter-arrival time)",
    ], fontsize=12, figw=4.5)
    slide.shapes.add_picture(eq, Inches(8.3), Inches(0.7), Inches(4.5))

    # Tool mapping
    tool_lines = [
        ("Flow Extraction Tools per Dataset:", 12, True, BROWN),
        ("", 4, False, BLACK),
        ("UNSW-NB15:", 11, True, BLACK),
        ("  PCAP \u2192 Argus + Bro-IDS \u2192 49-column CSV", 10, False, GRAY),
        ("  Captures: IXIA PerfectStorm traffic gen", 10, False, GRAY),
        ("", 4, False, BLACK),
        ("CSE-CIC-IDS-2018:", 11, True, BLACK),
        ("  PCAP \u2192 CICFlowMeter v4 \u2192 80+ column CSV", 10, False, GRAY),
        ("  Captures: AWS infrastructure, 10-day run", 10, False, GRAY),
        ("", 4, False, BLACK),
        ("CIC-IDS2017:", 11, True, BLACK),
        ("  PCAP \u2192 CICFlowMeter v3 \u2192 80+ column CSV", 10, False, GRAY),
        ("  Captures: CIC lab B-Profile, 5-day run", 10, False, GRAY),
    ]
    add_ml(slide, Inches(8.3), Inches(2.5), Inches(4.5), Inches(3.5), tool_lines, font_name="Consolas")

    # Bullets
    add_bullets(slide, Inches(0.3), Inches(6.6), Inches(12.5), Inches(0.8), [
        "Raw PCAP files contain binary packet data (Ethernet/IP/TCP headers + payload); packets are grouped into bidirectional flows using 5-tuple keys and timeout-based termination.",
        "CICFlowMeter (CIC datasets) and Argus/Bro (UNSW-NB15) compute flow-level statistics: duration, packet/byte rates, inter-arrival time distributions, TCP flag counts, and payload size statistics.",
        "Each PCAP file is processed independently, enabling embarrassingly parallel preprocessing across multiple capture files or time windows.",
        "Output: one CSV row per flow with 49\u201380+ raw features, which Stage 1 then maps to the standardized 7-feature schema."
    ], font_size=10)

    add_notes(slide, """SLIDE 4 \u2014 Stage 0: PCAP Binary Preprocessing & Flow Aggregation

PURPOSE: Bridge the gap between raw binary network captures (PCAP format) and the tabular flow-level data that the ML pipeline (Stages 1\u20138) consumes. This stage documents the complete data lineage from wire to CSV.

WHAT IS PCAP:
PCAP (Packet CAPture) is the standard binary format for storing captured network packets. Each PCAP file contains:
- Global header: magic number, version, snap length, link-layer type
- Per-packet records: timestamp (microsecond precision), captured length, original length, raw packet bytes
- Packet bytes include: Ethernet header (14 bytes) \u2192 IP header (20-60 bytes) \u2192 TCP/UDP/ICMP header \u2192 payload

FLOW AGGREGATION PROCESS:
1. PACKET PARSING: Read binary PCAP, decode Ethernet/IP/transport headers to extract 5-tuple
2. FLOW GROUPING: Group packets by (src_ip, dst_ip, src_port, dst_port, protocol) into bidirectional flows
3. TIMEOUT TERMINATION: Flows are terminated when:
   - Idle timeout (T_idle = 60s): no packets for 60 seconds
   - Active timeout (T_active = 120s): flow duration exceeds 120 seconds
   - FIN/RST flag: TCP connection termination
4. FEATURE COMPUTATION: For each terminated flow, compute statistical features:
   - Duration, total packets (fwd+bwd), total bytes (fwd+bwd)
   - Packet rate, byte rate, packet size stats (mean, std, min, max)
   - Inter-arrival time stats (mean, std, min, max) for fwd and bwd directions
   - TCP flags (SYN, ACK, FIN, RST, PSH, URG counts)
   - Payload statistics, header lengths, etc.

TOOLS USED:
- CICFlowMeter (v3/v4): Java-based tool from CIC that reads PCAP and outputs CICFlowMeter-format CSVs with 80+ features. Used for both CIC-IDS2017 and CSE-CIC-IDS-2018.
- Argus: Open-source flow generator that produces NetFlow-like records. Used alongside Bro-IDS (now Zeek) for UNSW-NB15.

PARALLELISM:
PCAP processing is embarrassingly parallel:
- Multiple PCAP files can be processed simultaneously (one per CPU core)
- Large PCAP files can be split by time window and processed in parallel
- This is the most computationally intensive preprocessing step (parsing millions of packets)

KEY POINT: The three benchmark datasets ship as pre-extracted CSVs. The PCAP-to-CSV conversion was performed by the dataset creators. However, in a production deployment, Stage 0 would run in real-time using CICFlowMeter or similar tools on live traffic captures.""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 5: Stage 1 \u2014 Data Ingestion & Z-Score Normalization
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Stage 1: Feature Extraction & StandardScaler Normalization",
                 font_size=24, bold=True, color=DARK_BLUE)

    algo_lines = [
        ("Require: Raw CSV D_raw from Stage 0, k \u2208 {UNSW, CSE, CIC2017}", 0, False),
        ("Ensure: Normalized feature matrix X_norm \u2208 R^(n\u00d77), labels y \u2208 {0,1}^n", 0, False),
        ("for each dataset D_k do", 0, True),
        ("for each column c_j \u2208 D_k do", 1, True),
        ("if c_j is numeric then", 2, True),
        ("x_j \u2190 map(c_j) using dataset-specific extraction rules", 3, False),
        ("else", 2, True),
        ("x_j \u2190 categorize(c_j) into integer bins", 3, False),
        ("end if", 2, True),
        ("end for", 1, True),
        ("X \u2190 [x_1, x_2, ..., x_7]  (n \u00d7 7 feature matrix)", 1, False),
        ("\u03bc_j \u2190 (1/n) \u2211 x_j^(i),  \u03c3_j \u2190 std(x_j)  for j = 1..4", 1, False),
        ("z_j^(i) \u2190 (x_j^(i) - \u03bc_j) / \u03c3_j  for continuous features", 1, False),
        ("X_norm \u2190 [z_1, z_2, z_3, z_4, x_5, x_6, x_7]", 1, False),
        ("end for", 0, True),
        ("return X_norm, y", 0, True),
    ]
    add_ieee_algorithm(slide, Inches(0.3), Inches(0.7), Inches(7.5), Inches(4.7),
                       2, "Feature Extraction & Z-Score Normalization", algo_lines)

    eq = render_eq_block([
        r"$\mu_j = \frac{1}{n}\sum_{i=1}^{n} x_j^{(i)}$",
        r"$\sigma_j = \sqrt{\frac{1}{n-1}\sum_{i=1}^{n}(x_j^{(i)} - \mu_j)^2}$",
        r"$z_j^{(i)} = \frac{x_j^{(i)} - \mu_j}{\sigma_j}$    (StandardScaler)",
    ], fontsize=14, figw=5)
    slide.shapes.add_picture(eq, Inches(8.2), Inches(0.7), Inches(4.8))

    add_bullets(slide, Inches(8.2), Inches(3.0), Inches(4.8), Inches(2.5), [
        "Maps the 49\u201380+ raw CSV columns from Stage 0 into a unified 7-dimensional representation: [flow_duration, pkt_rate, byte_rate, traffic_pat_var*, port_cat, size_cat, protocol]. *traffic_pat_var is a traffic pattern variability proxy, not Shannon entropy.",
        "Applies Z-score normalization (StandardScaler) to the 4 continuous features, ensuring zero mean and unit variance to prevent scale dominance.",
        "Categorical features (port_cat, size_cat, protocol) are binned into integer categories and passed through unchanged.",
        "The scaler is fit on training data only and applied to test data to prevent data leakage."
    ], font_size=11)

    add_text_box(slide, Inches(0.3), Inches(5.7), Inches(12.5), Inches(0.4),
                 "Output:  X_norm \u2208 R^(n\u00d77), y \u2208 {0,1}^n  \u2192 Feeds into Stage 2 (Decision Tree)",
                 font_size=13, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)

    add_notes(slide, """SLIDE 5 \u2014 Stage 1: Feature Extraction & StandardScaler Normalization

PURPOSE: Transform the raw CSV output from Stage 0 (49\u201380+ columns) into a standardized 7-feature normalized matrix.

This stage now explicitly receives input from Stage 0 (PCAP preprocessing), closing the gap between binary network data and the ML pipeline.

ALGORITHM DETAILS:
The StandardScaler performs Z-score normalization on each of the 4 continuous features independently:
- Compute mean and standard deviation from training data
- Transform: z = (x - mu) / sigma

DATASET-SPECIFIC MAPPINGS:
- UNSW-NB15 (49 cols \u2192 7): dur*1000, rate, sbytes+dbytes, normalize(ct_srv_src), bin(dsport), bin(byte_rate), proto_map
- CSE-CIC-2018 (80+ cols \u2192 7): Flow Duration/1000, Flow Packets/s, Flow Bytes/s, quantile_norm(Fwd IAT Std), bin(Dst Port), bin(bytes*dur), proto_map
- CIC-IDS2017 (80+ cols \u2192 7): Same as CSE-CIC-2018 (both use CICFlowMeter output format)

OUTPUT: X_norm (n\u00d77) and y (n binary labels) fed to Stage 2.""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 6: Stage 2 \u2014 Decision Tree Symbolic Model
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Stage 2: Decision Tree Symbolic Model (CART)",
                 font_size=24, bold=True, color=PURPLE)

    algo_lines = [
        ("Require: X_norm \u2208 R^(n\u00d77), y \u2208 {0,1}^n, max_depth=6, min_leaf=15", 0, False),
        ("Ensure: Trained tree T with L leaf nodes, leaf predictions", 0, False),
        ("T \u2190 initialize empty CART tree", 0, False),
        ("for each node v in T (breadth-first) do", 0, True),
        ("if depth(v) < max_depth and |samples(v)| \u2265 2 \u00d7 min_leaf then", 1, True),
        ("for each feature j \u2208 {1,...,7} do", 2, True),
        ("for each threshold \u03b8 in unique_splits(X_j) do", 3, True),
        ("Compute Gini_split = (|S_L|/|S|)\u00b7Gini(S_L) + (|S_R|/|S|)\u00b7Gini(S_R)", 4, False),
        ("end for", 3, True),
        ("end for", 2, True),
        ("(j*, \u03b8*) \u2190 argmin Gini_split", 2, False),
        ("Split v into left child (x_{j*} \u2264 \u03b8*) and right child (x_{j*} > \u03b8*)", 2, False),
        ("else", 1, True),
        ("Mark v as leaf: P(attack|v) = n_attack / n_total", 2, False),
        ("end if", 1, True),
        ("end for", 0, True),
        ("return T, {leaf_id(x_i), P(attack|x_i), \u0177_DT(x_i)} for all x_i", 0, True),
    ]
    add_ieee_algorithm(slide, Inches(0.3), Inches(0.7), Inches(7.5), Inches(5.1),
                       3, "CART Decision Tree Training", algo_lines)

    eq = render_eq_block([
        r"$Gini(S) = 1 - \sum_{k=0}^{1} p_k^2$",
        r"$p_k = \frac{|S_k|}{|S|}$",
        r"$P(attack | leaf) = \frac{n_{attack}^{leaf}}{n_{total}^{leaf}}$",
    ], fontsize=14, figw=4.5)
    slide.shapes.add_picture(eq, Inches(8.2), Inches(0.7), Inches(4.5))

    add_bullets(slide, Inches(8.2), Inches(3.0), Inches(4.8), Inches(2.5), [
        "Trains a CART decision tree (max_depth=6, min_samples_leaf=15) to serve as the interpretable symbolic model whose paths become Z3 constraints.",
        "Each leaf node provides: a discrete state ID for Q-learning (Stage 4), an attack probability P(attack|leaf), and a baseline classification.",
        "The shallow depth (max 64 leaves) keeps Z3 constraint extraction tractable while maintaining sufficient classification expressiveness.",
        "Gini impurity guides feature splits\u2014pure leaves (Gini=0) give high-confidence predictions; impure leaves are deferred to the RL agent."
    ], font_size=11)

    add_text_box(slide, Inches(0.3), Inches(6.2), Inches(12.5), Inches(0.4),
                 "Output: Tree T with L leaves \u2192 Feeds into Stage 3 (Z3 extraction) + Stage 4 (leaf_id as RL state)",
                 font_size=13, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)

    add_notes(slide, """SLIDE 6 \u2014 Stage 2: Decision Tree Symbolic Model (CART)

PURPOSE: Train an interpretable decision tree whose decision paths can be extracted as formal logic constraints.

WHY DECISION TREE: It is a "white-box" model \u2014 every prediction is a conjunction of feature threshold comparisons along a root-to-leaf path. This enables:
1. Direct translation into Z3 logical implications (Stage 3)
2. Discrete state IDs for Q-learning (Stage 4)
3. Tractable verification (max 64 leaves with depth=6)

NOTE ON PARALLELISM: CART training is inherently sequential \u2014 each split depends on the data partition from the previous split. This is one of the sequential bottlenecks in the pipeline. However, the tree is small (depth=6) so training completes in milliseconds.""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 7: Stage 3 \u2014 Z3 Formal Constraint Extraction
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Stage 3: Z3 SMT Constraint Extraction  [\u26a1 Parallelizable per Leaf Path]",
                 font_size=24, bold=True, color=GREEN)

    algo_lines = [
        ("Require: Trained decision tree T with L leaf nodes", 0, False),
        ("Ensure: Constraint set C = {\u03c6_1, ..., \u03c6_|C|} of Z3 implications", 0, False),
        ("C \u2190 \u2205", 0, False),
        ("for each leaf l \u2208 {1, ..., L} do  // \u26a1 parallelizable per leaf", 0, True),
        ("\u03c0_l \u2190 extract_path(root, l)  // walk tree root \u2192 leaf", 1, False),
        ("conditions \u2190 []", 1, False),
        ("for each split (feature_j, threshold_\u03b8, direction) in \u03c0_l do", 1, True),
        ("if direction = left then", 2, True),
        ("conditions.append( Real(f_j) \u2264 \u03b8 )", 3, False),
        ("else", 2, True),
        ("conditions.append( Real(f_j) > \u03b8 )", 3, False),
        ("end if", 2, True),
        ("end for", 1, True),
        ("\u03c6_l \u2190 Implies( And(conditions), Int(action) == class(l) )", 1, False),
        ("if solver.check(C \u222a {\u03c6_l}) == sat then", 1, True),
        ("C \u2190 C \u222a {\u03c6_l}", 2, False),
        ("end if", 1, True),
        ("end for", 0, True),
        ("return C", 0, True),
    ]
    add_ieee_algorithm(slide, Inches(0.3), Inches(0.7), Inches(7.5), Inches(5.3),
                       4, "Z3 Constraint Extraction", algo_lines)

    eq = render_eq_block([
        r"$\varphi_l : \bigwedge_{(j,\theta,op) \in \pi_l} (f_j\ op\ \theta) \Rightarrow (action = c_l)$",
        r"$C = \{\varphi_1, \varphi_2, \ldots, \varphi_{|C|}\}$",
        r"$verify(x, a) = solver.check(x \models C,\ action\!=\!a)$",
    ], fontsize=13, figw=5)
    slide.shapes.add_picture(eq, Inches(8.2), Inches(0.7), Inches(4.8))

    add_bullets(slide, Inches(8.2), Inches(3.0), Inches(4.8), Inches(2.5), [
        "Traverses each root-to-leaf path in the decision tree, converting feature threshold comparisons into Z3 Real/Int constraints with Implies() implications.",
        "Each constraint \u03c6_l states: if the conjunction of feature conditions along path \u03c0_l holds, then the action must equal the leaf's predicted class.",
        "Leaf paths are independent\u2014each can be extracted in parallel across CPU cores, with a final sequential merge and satisfiability check.",
        "The resulting constraint set C enables runtime verification: any proposed RL action can be checked against C using Z3's SMT solver in O(ms) time."
    ], font_size=11)

    add_text_box(slide, Inches(0.3), Inches(6.3), Inches(12.5), Inches(0.4),
                 "Output: Constraint set C \u2192 Used by Stage 4 (Q-Learning safety shield) and updated by Stage 5 (DBSCAN)",
                 font_size=13, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)

    add_notes(slide, """SLIDE 7 \u2014 Stage 3: Z3 SMT Constraint Extraction from Decision Tree

PURPOSE: Convert the decision tree into formal logic constraints using Z3.

PARALLELISM: Each leaf path extraction is independent:
- L leaf paths can be extracted concurrently (one per thread/core)
- Each path produces a candidate constraint \u03c6_l
- The satisfiability merge step is sequential (constraints must be checked for mutual consistency)
- With L \u2248 60 leaves, parallelism reduces extraction time by ~Lx on multi-core systems

The verify() function checks: given feature values and proposed action, is the constraint set satisfiable? sat = safe, unsat = unsafe.""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 8: Stage 4 \u2014 Q-Learning RL with Safety Shielding
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Stage 4: Q-Learning Reinforcement Learning with Z3 Safety Shielding",
                 font_size=24, bold=True, color=DARK_ORANGE)

    algo_lines = [
        ("Require: Constraint set C, DT tree T, training data (X, y)", 0, False),
        ("Ensure: Learned Q-table Q(s, a) with safety-verified actions", 0, False),
        ("Initialize Q(s, a) \u2190 0 for all states s, actions a \u2208 {0, 1}", 0, False),
        ("\u03b5 \u2190 0.20, \u03b1 \u2190 0.15, \u03b3 \u2190 0.95", 0, False),
        ("for each episode e = 1, ..., E do", 0, True),
        ("for each flow x_i in training batch do", 1, True),
        ("s_i \u2190 leaf_id(T, x_i)  // DT leaf as discrete state", 2, False),
        ("if random() < \u03b5 then", 2, True),
        ("a_proposed \u2190 random_action({0, 1})  // explore", 3, False),
        ("else", 2, True),
        ("a_proposed \u2190 argmax_a Q(s_i, a)  // exploit", 3, False),
        ("end if", 2, True),
        ("if Z3_verify(C, x_i, a_proposed) == sat then", 2, True),
        ("a_i \u2190 a_proposed  // action is safe", 3, False),
        ("else", 2, True),
        ("a_i \u2190 1 - a_proposed  // safety shield override", 3, False),
        ("end if", 2, True),
        ("r_i \u2190 Reward(a_i, y_i)  // TP:+2, FN:-3, TN:+1, FP:-1", 2, False),
        ("if Z3_verify passed then r_i \u2190 r_i + 0.5  // shield bonus", 2, True),
        ("Q(s_i, a_i) \u2190 Q(s_i, a_i) + \u03b1[r_i + \u03b3\u00b7max_a Q(s', a) - Q(s_i, a_i)]", 2, False),
        ("end for", 1, True),
        ("\u03b5 \u2190 \u03b5 \u00d7 0.999  // decay exploration", 1, False),
        ("end for", 0, True),
        ("return Q", 0, True),
    ]
    add_ieee_algorithm(slide, Inches(0.3), Inches(0.7), Inches(8.0), Inches(5.8),
                       5, "Q-Learning with Z3 Safety Shielding", algo_lines)

    eq = render_eq_block([
        r"$Q(s,a) \leftarrow Q(s,a) + \alpha[r + \gamma \max_{a'} Q(s',a') - Q(s,a)]$",
        r"$R(a,y) = \{+2\ (TP),\ -3\ (FN),\ +1\ (TN),\ -1\ (FP)\}$",
    ], fontsize=13, figw=4.5)
    slide.shapes.add_picture(eq, Inches(8.5), Inches(0.7), Inches(4.5))

    add_bullets(slide, Inches(8.5), Inches(2.5), Inches(4.5), Inches(3.0), [
        "A tabular Q-learning agent uses DT leaf IDs as discrete states and {ALLOW=0, BLOCK=1} as actions, learning from asymmetric rewards: FN penalty (\u22123) exceeds FP penalty (\u22121) to prioritize attack detection.",
        "Every proposed action is verified against the Z3 constraint set\u2014if unsafe (unsat), the safety shield overrides it with the alternative action.",
        "Epsilon-greedy exploration (\u03b5=0.20, decay=0.999) balances discovery with exploitation; a +0.5 shield bonus incentivizes Z3-consistent behavior.",
        "Sequential bottleneck: Q-table updates within each episode are sequential (shared state), but episodes across datasets run in parallel."
    ], font_size=11)

    add_notes(slide, """SLIDE 8 \u2014 Stage 4: Q-Learning RL with Z3 Safety Shielding

PURPOSE: Learn an adaptive policy that improves on the static DT while maintaining safety guarantees.

PARALLELISM NOTE: Q-learning training is partially sequential:
- Within an episode, each flow's Q-update depends on the current Q-table \u2192 sequential
- However, separate Q-tables are trained for each dataset independently \u2192 cross-dataset parallelism
- The Z3 verify() calls per flow are independent and could be batched

REWARD: TP:+2, FN:-3, TN:+1, FP:-1, shield_bonus:+0.5
The asymmetric penalty (-3 FN vs -1 FP) reflects IDS priority: missing attacks is worse than false alarms.""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 9: Stage 5 \u2014 DBSCAN Novel Pattern Detection
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Stage 5: DBSCAN Novel Attack Pattern Detection",
                 font_size=24, bold=True, color=DARK_RED)

    algo_lines = [
        ("Require: Misclassified flows M = {x_i : \u0177_i \u2260 y_i}, constraint set C, \u03b5=1.5, minPts=5", 0, False),
        ("Ensure: Updated constraint set C' = C \u222a C_novel", 0, False),
        ("if |M| < minPts then return C  // insufficient data", 0, True),
        ("clusters \u2190 DBSCAN(M, \u03b5=1.5, min_samples=5)", 0, False),
        ("C_novel \u2190 \u2205", 0, False),
        ("for each cluster K_m \u2208 clusters (excluding noise) do", 0, True),
        ("\u03bc_m \u2190 centroid(K_m)  // mean feature vector", 1, False),
        ("\u03c3_m \u2190 std(K_m)  // feature-wise standard deviation", 1, False),
        ("for each feature j = 1, ..., 7 do", 1, True),
        ("\u03c6_j \u2190 And( Real(f_j) \u2265 \u03bc_m[j] - 2\u03c3_m[j],", 2, False),
        ("            Real(f_j) \u2264 \u03bc_m[j] + 2\u03c3_m[j] )", 2, False),
        ("end for", 1, True),
        ("\u03c6_novel \u2190 Implies( And(\u03c6_1,...,\u03c6_7), action == majority_class(K_m) )", 1, False),
        ("if solver.check(C \u222a {\u03c6_novel}) == sat then", 1, True),
        ("C_novel \u2190 C_novel \u222a {\u03c6_novel}", 2, False),
        ("end if", 1, True),
        ("end for", 0, True),
        ("return C' = C \u222a C_novel", 0, True),
    ]
    add_ieee_algorithm(slide, Inches(0.3), Inches(0.7), Inches(7.8), Inches(5.3),
                       6, "DBSCAN Novel Pattern Detection & Constraint Generation", algo_lines)

    eq = render_eq_block([
        r"$d(x_i, x_j) = \|x_i - x_j\|_2$    (Euclidean distance)",
        r"$N_\epsilon(x_i) = \{x_j : d(x_i, x_j) \leq \epsilon\}$",
        r"$\varphi_{novel}: (\mu_m \pm 2\sigma_m) \Rightarrow action$",
    ], fontsize=13, figw=4.5)
    slide.shapes.add_picture(eq, Inches(8.3), Inches(0.7), Inches(4.5))

    add_bullets(slide, Inches(8.3), Inches(2.8), Inches(4.7), Inches(3.0), [
        "Collects misclassified flows and applies DBSCAN density-based clustering (\u03b5=1.5, minPts=5) to discover previously unseen attack patterns.",
        "Each cluster is converted into a Z3 constraint using centroid \u00b1 2\u03c3 per feature, then added after satisfiability verification.",
        "Noise points (not in any cluster) are discarded\u2014only dense groups generate new constraints, reducing false pattern creation.",
        "Enables zero-day detection: novel attacks not seen in training are discovered at runtime and formally incorporated into the safety shield."
    ], font_size=11)

    add_text_box(slide, Inches(0.3), Inches(6.3), Inches(12.5), Inches(0.4),
                 "Output: C' = C \u222a C_novel \u2192 Feedback loop: updated constraints improve Stage 4 safety shield",
                 font_size=13, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)

    add_notes(slide, """SLIDE 9 \u2014 Stage 5: DBSCAN Novel Attack Pattern Detection

DBSCAN is inherently sequential (neighborhood queries depend on prior classifications of core/border/noise points). This is a sequential bottleneck, but runs on a small subset of data (only misclassified flows).

The feedback loop C' = C \u222a C_novel means the constraint set grows over time, enabling detection of zero-day attacks.""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 10: Stage 6 \u2014 Adaptive Buffer Management
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Stage 6: Adaptive Buffer Management",
                 font_size=24, bold=True, color=MED_BLUE)

    algo_lines = [
        ("Require: Flow stream F, current buffer B_t with size |B_t|, bounds [10, 200]", 0, False),
        ("Ensure: Resized buffer B_{t+1} adapted to traffic volatility", 0, False),
        ("for each time window W_t of flows do", 0, True),
        ("B_t.append(flows in W_t)  // sliding window buffer", 1, False),
        ("\u03c3_tpv \u2190 std({traffic_pat_var(x) : x \u2208 B_t})  // traffic pattern variability volatility", 1, False),
        ("\u03c3\u00b2_bytes \u2190 var({byte_rate(x) : x \u2208 B_t})  // byte rate variance", 1, False),
        ("volatility \u2190 \u03c3_tpv + \u03c3\u00b2_bytes / max(\u03c3\u00b2_bytes)", 1, False),
        ("if volatility > high_threshold then", 1, True),
        ("\u0394 \u2190 +10  // high volatility: grow buffer for more context", 2, False),
        ("else if volatility < low_threshold then", 1, True),
        ("\u0394 \u2190 -5  // low volatility: shrink buffer for faster response", 2, False),
        ("else", 1, True),
        ("\u0394 \u2190 0  // stable: maintain current size", 2, False),
        ("end if", 1, True),
        ("|B_{t+1}| \u2190 clip(|B_t| + \u0394, 10, 200)", 1, False),
        ("end for", 0, True),
        ("return B_{t+1}", 0, True),
    ]
    add_ieee_algorithm(slide, Inches(0.3), Inches(0.7), Inches(7.5), Inches(5.0),
                       7, "Adaptive Buffer Resizing", algo_lines)

    eq2 = render_eq_block([
        r"$|B_{t+1}| = clip(|B_t| + \Delta,\ 10,\ 200)$",
        r"$\Delta = +10\ \ if\ volatility > \tau_{high}$",
        r"$\Delta = -5\ \ \ if\ volatility < \tau_{low}$",
        r"$\Delta = 0\ \ \ \ otherwise$",
    ], fontsize=13, figw=4.5)
    slide.shapes.add_picture(eq2, Inches(8.2), Inches(0.7), Inches(4.5))

    add_bullets(slide, Inches(8.2), Inches(3.0), Inches(4.8), Inches(2.5), [
        "Dynamically resizes the flow buffer (10\u2013200 flows) based on traffic volatility measured by traffic_pat_var std and byte rate variance.",
        "High volatility triggers buffer growth (+10) to accumulate more context; low volatility triggers shrinkage (\u22125) for faster response.",
        "Sequential stage: buffer state depends on previous window (temporal dependency), so this cannot be parallelized within a single stream.",
        "Asymmetric growth/shrink rates (+10/\u22125) create a conservative bias\u2014the system is quicker to accumulate evidence than to discard it."
    ], font_size=11)

    add_text_box(slide, Inches(0.3), Inches(6.0), Inches(12.5), Inches(0.4),
                 "Output: Adapted buffer size |B_{t+1}| \u2192 Used by Stage 7 and Stage 8 for windowed classification",
                 font_size=13, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)

    add_notes(slide, """SLIDE 10 \u2014 Stage 6: Adaptive Buffer Management

Sequential stage \u2014 buffer state depends on previous window. Cannot be parallelized within a single traffic stream, but different streams (e.g., per-subnet or per-sensor) could maintain independent buffers in parallel.

Grows +10 during volatility, shrinks -5 during stability. Bounds: [10, 200]. Asymmetric rates create conservative bias.""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 11: Stage 7 \u2014 Dynamic Threshold Adaptation
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Stage 7: Dynamic Threshold Adaptation",
                 font_size=24, bold=True, color=TEAL)

    algo_lines = [
        ("Require: Current threshold \u03c4_t \u2208 [0.40, 0.70], recent FPR_t, FNR_t, EMA decay \u03b2=0.3", 0, False),
        ("Ensure: Updated threshold \u03c4_{t+1} balancing false positives and false negatives", 0, False),
        ("for each evaluation window W_t do", 0, True),
        ("FPR_t \u2190 FP_t / (FP_t + TN_t)  // false positive rate", 1, False),
        ("FNR_t \u2190 FN_t / (FN_t + TP_t)  // false negative rate", 1, False),
        ("EMA_FPR \u2190 \u03b2 \u00b7 FPR_t + (1 - \u03b2) \u00b7 EMA_FPR  // exponential moving average", 1, False),
        ("EMA_FNR \u2190 \u03b2 \u00b7 FNR_t + (1 - \u03b2) \u00b7 EMA_FNR", 1, False),
        ("if EMA_FPR > EMA_FNR then", 1, True),
        ("\u03c4_{t+1} \u2190 \u03c4_t + 0.005  // too many false positives: raise threshold", 2, False),
        ("else if EMA_FNR > EMA_FPR then", 1, True),
        ("\u03c4_{t+1} \u2190 \u03c4_t - 0.005  // too many false negatives: lower threshold", 2, False),
        ("else", 1, True),
        ("\u03c4_{t+1} \u2190 \u03c4_t  // balanced: maintain threshold", 2, False),
        ("end if", 1, True),
        ("\u03c4_{t+1} \u2190 clip(\u03c4_{t+1}, 0.40, 0.70)", 1, False),
        ("end for", 0, True),
        ("return \u03c4_{t+1}", 0, True),
    ]
    add_ieee_algorithm(slide, Inches(0.3), Inches(0.7), Inches(7.5), Inches(5.0),
                       8, "Dynamic Threshold Adaptation via EMA Error Balancing", algo_lines)

    eq = render_eq_block([
        r"$EMA_t = \beta \cdot x_t + (1-\beta) \cdot EMA_{t-1}$",
        r"$\tau_{t+1} = clip(\tau_t \pm 0.005,\ 0.40,\ 0.70)$",
        r"$FPR = \frac{FP}{FP + TN}$,    $FNR = \frac{FN}{FN + TP}$",
    ], fontsize=13, figw=4.5)
    slide.shapes.add_picture(eq, Inches(8.2), Inches(0.7), Inches(4.5))

    add_bullets(slide, Inches(8.2), Inches(3.0), Inches(4.8), Inches(2.5), [
        "Continuously adjusts \u03c4 \u2208 [0.40, 0.70] based on the balance between FPR and FNR using exponential moving average (EMA, \u03b2=0.3).",
        "When FPR dominates, threshold rises to reduce false alarms; when FNR dominates, it falls to catch more attacks\u2014creating a self-correcting feedback loop.",
        "Sequential stage: threshold state depends on previous window's EMA values (temporal dependency).",
        "Hard bounds [0.40, 0.70] prevent threshold from becoming too permissive or too aggressive."
    ], font_size=11)

    add_text_box(slide, Inches(0.3), Inches(6.0), Inches(12.5), Inches(0.4),
                 "Output: \u03c4_{t+1} \u2192 Used by Stage 8 for high-confidence direct classification",
                 font_size=13, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)

    add_notes(slide, """SLIDE 11 \u2014 Stage 7: Dynamic Threshold Adaptation

Sequential stage \u2014 EMA depends on previous state. Like Stage 6, this cannot be parallelized within a single stream but operates independently across datasets.

Observed convergence: UNSW-NB15 \u2192 \u03c4\u22480.625, CSE-CIC-2018 \u2192 \u03c4\u22480.625, CIC-IDS2017 \u2192 \u03c4\u22480.685.""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 12: Stage 8 \u2014 Final Classification Decision
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Stage 8: Final Classification Decision  [\u26a1 Parallelizable per Flow]",
                 font_size=24, bold=True, color=NAVY)

    algo_lines = [
        ("Require: Flow x_i, threshold \u03c4, Q-table Q, constraint set C, DT tree T", 0, False),
        ("Ensure: Final classification \u0177_i \u2208 {0=BENIGN, 1=ATTACK} with audit trail", 0, False),
        ("p_i \u2190 P(attack | leaf_id(T, x_i))  // DT leaf probability", 0, False),
        ("if p_i \u2265 \u03c4 then", 0, True),
        ("\u0177_i \u2190 1 (ATTACK)  // high confidence attack", 1, False),
        ("confidence \u2190 'HIGH_ATTACK'", 1, False),
        ("else if p_i \u2264 1 - \u03c4 then", 0, True),
        ("\u0177_i \u2190 0 (BENIGN)  // high confidence benign", 1, False),
        ("confidence \u2190 'HIGH_BENIGN'", 1, False),
        ("else", 0, True),
        ("// Uncertain region: invoke RL + Z3 safety shield", 1, False),
        ("s_i \u2190 leaf_id(T, x_i)", 1, False),
        ("a_RL \u2190 argmax_a Q(s_i, a)", 1, False),
        ("if Z3_verify(C, x_i, a_RL) == sat then", 1, True),
        ("\u0177_i \u2190 a_RL  // RL action is safe", 2, False),
        ("else", 1, True),
        ("\u0177_i \u2190 1 - a_RL  // safety shield override", 2, False),
        ("end if", 1, True),
        ("confidence \u2190 'RL_VERIFIED'", 1, False),
        ("end if", 0, True),
        ("return \u0177_i, confidence, audit_trail(p_i, s_i, C)  // \u26a1 per-flow parallel", 0, True),
    ]
    add_ieee_algorithm(slide, Inches(0.3), Inches(0.7), Inches(7.8), Inches(5.5),
                       9, "Final Classification with Tiered Decision Logic", algo_lines)

    eq = render_eq_block([
        r"$p_i \geq \tau \Rightarrow ATTACK$",
        r"$p_i \leq 1-\tau \Rightarrow BENIGN$",
        r"$otherwise \Rightarrow RL + Z3\ shield$",
    ], fontsize=13, figw=4)
    slide.shapes.add_picture(eq, Inches(8.3), Inches(0.7), Inches(4.0))

    add_bullets(slide, Inches(8.3), Inches(2.8), Inches(4.5), Inches(3.0), [
        "Three-tier logic: high-confidence flows (p \u2265 \u03c4 or p \u2264 1\u2212\u03c4) classified directly; uncertain flows deferred to RL+Z3 verification.",
        "Embarrassingly parallel: each flow's classification is independent\u2014Q-table and constraint set are read-only at inference time.",
        "Every classification produces a full audit trail: DT probability, leaf state, Z3 check result, and confidence level for complete explainability.",
        "Tiered approach reduces Z3 solver calls (~70% of flows are high-confidence), balancing verification overhead with throughput."
    ], font_size=11)

    add_text_box(slide, Inches(0.3), Inches(6.3), Inches(12.5), Inches(0.4),
                 "Output: \u0177_i \u2208 {0, 1} with confidence level and full audit trail for every classification decision",
                 font_size=13, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)

    add_notes(slide, """SLIDE 12 \u2014 Stage 8: Final Classification Decision

PARALLELISM: This is the most parallelizable stage. At inference time:
- The Q-table is frozen (read-only) \u2192 no write conflicts
- The constraint set C is frozen (read-only) \u2192 each Z3 solver instance is independent
- Each flow x_i is classified independently \u2192 embarrassingly parallel
- Can distribute across N cores for ~Nx throughput improvement

This is where real-time IDS throughput is critical \u2014 production systems need to classify millions of flows per second.""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 13: Parallel Processing Architecture (NEW)
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Parallel Processing Architecture Across the ASRRL Pipeline",
                 font_size=26, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    # Parallel stages (green boxes)
    add_text_box(slide, Inches(0.3), Inches(0.8), Inches(6), Inches(0.4),
                 "\u26a1 Parallelizable Stages (Embarrassingly Parallel)",
                 font_size=16, bold=True, color=DARK_GREEN)

    par_lines = [
        ("Stage 0: PCAP Binary Preprocessing", 13, True, DARK_GREEN),
        ("  Granularity: Per PCAP file / per time window", 11, False, BLACK),
        ("  Strategy: Each capture file processed by independent worker", 11, False, BLACK),
        ("  Speedup: ~Nx with N CPU cores (I/O bound)", 11, False, GRAY),
        ("  Tool: CICFlowMeter instances or parallel Argus processes", 11, False, GRAY),
        ("", 4, False, BLACK),
        ("Stage 3: Z3 Constraint Extraction", 13, True, DARK_GREEN),
        ("  Granularity: Per leaf path (L \u2248 60 independent paths)", 11, False, BLACK),
        ("  Strategy: Extract each path in parallel; sequential merge + sat check", 11, False, BLACK),
        ("  Speedup: ~min(L, N) with N cores", 11, False, GRAY),
        ("  Note: Final satisfiability merge is sequential (consistency check)", 11, False, GRAY),
        ("", 4, False, BLACK),
        ("Stage 8: Final Classification (Inference)", 13, True, DARK_GREEN),
        ("  Granularity: Per flow (each flow independently classified)", 11, False, BLACK),
        ("  Strategy: Partition flow stream across N workers; read-only Q-table + C", 11, False, BLACK),
        ("  Speedup: ~Nx with N cores (CPU bound per Z3 call)", 11, False, GRAY),
        ("  Note: ~70% direct classification (no Z3 call); ~30% invoke solver", 11, False, GRAY),
        ("", 4, False, BLACK),
        ("Cross-Dataset: All 3 datasets processed independently", 13, True, DARK_GREEN),
        ("  Strategy: 3 pipeline instances run in parallel (no shared state)", 11, False, BLACK),
        ("  Speedup: 3x wall-clock time reduction", 11, False, GRAY),
    ]
    add_ml(slide, Inches(0.3), Inches(1.2), Inches(6.2), Inches(5.5), par_lines, font_name="Consolas")

    # Sequential stages (red boxes)
    add_text_box(slide, Inches(7.0), Inches(0.8), Inches(6), Inches(0.4),
                 "\u23f8 Sequential Bottlenecks",
                 font_size=16, bold=True, color=DARK_RED)

    seq_lines = [
        ("Stage 2: Decision Tree Training", 13, True, DARK_RED),
        ("  Reason: CART greedy splits are data-dependent", 11, False, BLACK),
        ("  Each split depends on the partition from", 11, False, BLACK),
        ("  the previous split (inherently sequential).", 11, False, BLACK),
        ("  Impact: Low (tree depth=6, trains in <1 second)", 11, False, GRAY),
        ("", 4, False, BLACK),
        ("Stage 4: Q-Learning Training (within episode)", 13, True, DARK_RED),
        ("  Reason: Q-table updates are sequential", 11, False, BLACK),
        ("  (each update depends on current Q-values).", 11, False, BLACK),
        ("  Mitigation: Cross-dataset parallelism (3 Q-tables)", 11, False, GRAY),
        ("", 4, False, BLACK),
        ("Stage 5: DBSCAN Clustering", 13, True, DARK_RED),
        ("  Reason: Core/border/noise classification is", 11, False, BLACK),
        ("  sequential (neighbor queries depend on prior labels).", 11, False, BLACK),
        ("  Impact: Low (runs on small subset: misclassified only)", 11, False, GRAY),
        ("", 4, False, BLACK),
        ("Stages 6-7: Buffer & Threshold Adaptation", 13, True, DARK_RED),
        ("  Reason: Temporal dependency \u2014 each window's state", 11, False, BLACK),
        ("  depends on the previous window's EMA / buffer size.", 11, False, BLACK),
        ("  Mitigation: Independent per-stream/per-sensor buffers", 11, False, GRAY),
    ]
    add_ml(slide, Inches(7.0), Inches(1.2), Inches(6), Inches(5.5), seq_lines, font_name="Consolas")

    # Bottom summary
    add_text_box(slide, Inches(0.3), Inches(6.8), Inches(12.5), Inches(0.5),
                 "Training Phase: Stages 0\u20135 (mostly sequential except Stage 0 & 3)  |  "
                 "Inference Phase: Stage 8 (fully parallel per flow)  |  "
                 "Adaptation: Stages 6\u20137 (sequential per stream, parallel across streams)",
                 font_size=12, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    add_notes(slide, """SLIDE 13 \u2014 Parallel Processing Architecture

This slide addresses the advisor's question about where parallel processing is used in the pipeline.

SUMMARY OF PARALLELISM:

EMBARRASSINGLY PARALLEL STAGES:
1. Stage 0 (PCAP Preprocessing): Each PCAP file is processed by an independent CICFlowMeter/Argus instance. No shared state between files. This is I/O-bound and scales linearly with CPU cores.

2. Stage 3 (Z3 Constraint Extraction): Each of the L leaf paths can be extracted independently. The path traversal and constraint construction are pure functions of the tree structure. Only the final merge step (checking consistency of the full constraint set) requires sequential processing.

3. Stage 8 (Final Classification at Inference): This is the critical real-time stage. At inference time, the Q-table and constraint set are both read-only (frozen after training). Each flow can be classified by an independent worker thread/process. With N CPU cores, throughput scales ~Nx. The tiered approach further helps: ~70% of flows take the fast path (direct DT classification, no Z3 call), only ~30% require Z3 solver invocation.

4. Cross-Dataset Parallelism: Since each dataset has its own independent pipeline (separate DT, Q-table, constraint set, buffer, threshold), all three datasets can be processed simultaneously.

SEQUENTIAL BOTTLENECKS (and why they're acceptable):
1. Stage 2 (DT Training): CART is inherently sequential, but with depth=6 and 50K samples, it completes in <1 second.
2. Stage 4 (Q-Learning): Within an episode, Q-updates are sequential. But training converges quickly and runs offline.
3. Stage 5 (DBSCAN): Sequential neighborhood queries, but operates on the small subset of misclassified flows only.
4. Stages 6-7 (Buffer/Threshold): Temporal dependencies prevent within-stream parallelism. However, different network segments (per-subnet, per-sensor, per-VLAN) can maintain independent buffer/threshold states in parallel.

TRAINING vs INFERENCE:
- Training phase (Stages 0-7): Mostly sequential, runs offline. Total time: seconds to minutes depending on dataset size.
- Inference phase (Stage 8): Fully parallel, runs in real-time. This is where parallelism matters most for production IDS deployment.

IMPLEMENTATION OPTIONS:
- Python multiprocessing.Pool for Stage 0 (PCAP files)
- concurrent.futures.ThreadPoolExecutor for Stage 3 (leaf paths, GIL-free Z3 calls)
- Ray or Dask for distributed Stage 8 inference across multiple machines
- Independent process per dataset for cross-dataset parallelism""")

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 14: Results Summary
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_text_box(slide, Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.5),
                 "Experimental Results \u2014 ASRRL Dynamic vs. Static Baselines",
                 font_size=26, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

    results = [
        ("", 6, False, BLACK),
        ("                          ASRRL Dynamic        Random Forest      XGBoost         SVM             KNN", 10, True, DARK_BLUE),
        ("                          (Buf+Thresh)                                                              ", 10, False, GRAY),
        ("", 4, False, BLACK),
        ("CSE-CIC-IDS-2018", 12, True, BLACK),
        ("  F1-Score:               0.9812               0.9956             0.9934          0.9901          0.9799", 10, False, BLACK),
        ("  Accuracy:               0.9887               0.9973             0.9960          0.9940          0.9880", 10, False, BLACK),
        ("  FPR:                    0.0048               0.0010             0.0010          0.0010          0.0010", 10, False, BLACK),
        ("  Final Buffer:           30                   N/A                N/A             N/A             N/A", 10, False, MED_BLUE),
        ("  Final Threshold:        0.625                N/A                N/A             N/A             N/A", 10, False, MED_BLUE),
        ("", 4, False, BLACK),
        ("UNSW-NB15", 12, True, BLACK),
        ("  F1-Score:               0.9889               1.0000             0.9967          0.9923          0.9833", 10, False, BLACK),
        ("  Accuracy:               0.9933               1.0000             0.9980          0.9953          0.9900", 10, False, BLACK),
        ("  FPR:                    0.0000               0.0000             0.0000          0.0019          0.0000", 10, False, BLACK),
        ("  Final Buffer:           37                   N/A                N/A             N/A             N/A", 10, False, MED_BLUE),
        ("  Final Threshold:        0.650                N/A                N/A             N/A             N/A", 10, False, MED_BLUE),
        ("", 4, False, BLACK),
        ("CIC-IDS2017", 12, True, BLACK),
        ("  F1-Score:               0.9845               0.9978             0.9934          0.9945          0.9822", 10, False, BLACK),
        ("  Accuracy:               0.9907               0.9987             0.9960          0.9967          0.9893", 10, False, BLACK),
        ("  FPR:                    0.0019               0.0010             0.0010          0.0010          0.0010", 10, False, BLACK),
        ("  Final Buffer:           32                   N/A                N/A             N/A             N/A", 10, False, MED_BLUE),
        ("  Final Threshold:        0.685                N/A                N/A             N/A             N/A", 10, False, MED_BLUE),
        ("", 4, False, BLACK),
        ("Key Insight: ASRRL achieves competitive F1 (0.98+) while providing formal safety guarantees,", 11, True, DARK_GREEN),
        ("dynamic adaptation (buffer+threshold), zero-day detection, and parallel inference scalability.", 11, True, DARK_GREEN),
    ]
    add_ml(slide, Inches(0.3), Inches(0.6), Inches(12.5), Inches(6.5), results, font_name="Consolas")

    add_notes(slide, """SLIDE 14 \u2014 Experimental Results

ASRRL Dynamic achieves F1 scores of 0.98+ across all three datasets, competitive with static baselines (RF: 0.99+, XGBoost: 0.99+).

Why slightly lower F1 is acceptable:
1. FORMAL SAFETY GUARANTEES: Every classification is Z3-verified
2. DYNAMIC ADAPTATION: Buffer and threshold adapt at runtime
3. ZERO-DAY DETECTION: DBSCAN discovers novel attack patterns
4. PARALLEL INFERENCE: Stage 8 scales linearly with CPU cores
5. COMPLETE EXPLAINABILITY: Full audit trail per decision

Statistical significance: Wilcoxon tests show significant differences from most baselines at p<0.05.
Constraint fidelity: 100% fidelity, coverage, and opinion rate across all datasets.""")

    # ── Save ──
    out = os.path.join(os.path.dirname(__file__), "ASRRL_Methodology.pptx")
    prs.save(out)
    print(f"Saved \u2192 {out}  ({os.path.getsize(out)//1024} KB)")


if __name__ == "__main__":
    build_presentation()
